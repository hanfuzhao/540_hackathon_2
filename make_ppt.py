from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Helpers ──────────────────────────────────────────────────────────────────

W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

BG_DARK   = RGBColor(0x0f, 0x0f, 0x1a)
PURPLE    = RGBColor(0xc0, 0x84, 0xfc)
GREEN     = RGBColor(0x4a, 0xde, 0x80)
AMBER     = RGBColor(0xfb, 0xbf, 0x24)
RED       = RGBColor(0xf8, 0x71, 0x71)
BLUE      = RGBColor(0x96, 0xc8, 0xfa)
WHITE     = RGBColor(0xf0, 0xf0, 0xf0)
DIMWHITE  = RGBColor(0xb0, 0xb0, 0xc0)

SLIDE_BGS = [
    RGBColor(0x2d, 0x0f, 0x5e),
    RGBColor(0x12, 0x24, 0x48),
    RGBColor(0x1a, 0x3d, 0x24),
    RGBColor(0x3d, 0x26, 0x00),
    RGBColor(0x2a, 0x0a, 0x4a),
]


def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=20, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_box(slide, left, top, width, height,
            fill_color=RGBColor(0x20, 0x20, 0x35),
            line_color=RGBColor(0x50, 0x50, 0x80),
            radius=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(0.75)
    return shape


# ── Build presentation ────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank

# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title & Problem
# ════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
set_bg(s1, SLIDE_BGS[0])

add_textbox(s1, "AIPI 540 · Mini Hackathon 2 · 2026",
            Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.5),
            font_size=12, color=DIMWHITE, align=PP_ALIGN.CENTER)

add_textbox(s1, "🛡️  Gaslight Guard",
            Inches(0.8), Inches(1.0), Inches(11.7), Inches(1.2),
            font_size=54, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

add_textbox(s1,
            "Can machines recognize veiled verbal abuse?\n"
            "Standard sentiment classifiers read gaslighting as positive.\n"
            "We built a system that doesn't.",
            Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.4),
            font_size=18, color=DIMWHITE, align=PP_ALIGN.CENTER)

# Label chips row (4 boxes)
chips = [
    ("Gaslighting",        RED),
    ("Passive-Aggressive", AMBER),
    ("Sarcastic",          PURPLE),
    ("Sincere",            GREEN),
]
chip_w = Inches(2.2)
chip_h = Inches(0.42)
gap    = Inches(0.3)
total_w = len(chips) * chip_w + (len(chips)-1) * gap
start_x = (W - total_w) / 2
for i, (label, color) in enumerate(chips):
    x = start_x + i * (chip_w + gap)
    b = add_box(s1, x, Inches(3.85), chip_w, chip_h,
                fill_color=RGBColor(0x22, 0x10, 0x35),
                line_color=color)
    add_textbox(s1, label, x, Inches(3.85), chip_w, chip_h,
                font_size=13, bold=True, color=color, align=PP_ALIGN.CENTER)

# Stats row
stats = [("4", "toxic-speech labels"), ("2", "evaluation axes"), ("5", "message context window")]
stat_w = Inches(3.0)
sx = Inches(1.2)
for val, lbl in stats:
    add_textbox(s1, val,  sx, Inches(4.65), stat_w, Inches(0.7),
                font_size=40, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    add_textbox(s1, lbl, sx, Inches(5.35), stat_w, Inches(0.45),
                font_size=12, color=DIMWHITE, align=PP_ALIGN.CENTER)
    sx += Inches(3.7)

add_textbox(s1, "01 / 05", Inches(11.8), Inches(0.15), Inches(1.3), Inches(0.35),
            font_size=10, color=DIMWHITE, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — NLP Approach
# ════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
set_bg(s2, SLIDE_BGS[1])

add_textbox(s2, "NLP Approach",
            Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
            font_size=36, bold=True, color=WHITE)
add_textbox(s2, "02 / 05", Inches(11.8), Inches(0.15), Inches(1.3), Inches(0.35),
            font_size=10, color=DIMWHITE, align=PP_ALIGN.RIGHT)

# Flow diagram — 5 boxes + arrows
flow_items = [
    ("Input",  "Message or\nconversation\n(≤5 turns)",          BLUE),
    ("Model",  "Fine-tuned\nRoBERTa\n+ zero-shot fallback",     PURPLE),
    ("Rules",  "Regex pattern\nlibrary\n(GL · PA · sarcasm)",   GREEN),
    ("Blend",  "85% model\n+ 15% rules\n(zero-shot: 70/30)",    AMBER),
    ("Output", "Label + severity\n+ coping\nstrategies",        RED),
]
box_w   = Inches(2.1)
box_h   = Inches(1.45)
flow_y  = Inches(1.35)
arrow_w = Inches(0.35)
total   = len(flow_items) * box_w + (len(flow_items)-1) * arrow_w
fx      = (W - total) / 2

for i, (title, body, color) in enumerate(flow_items):
    add_box(s2, fx, flow_y, box_w, box_h,
            fill_color=RGBColor(0x10, 0x1a, 0x2e),
            line_color=color)
    add_textbox(s2, title, fx, flow_y + Inches(0.08), box_w, Inches(0.35),
                font_size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(s2, body, fx, flow_y + Inches(0.42), box_w, Inches(0.95),
                font_size=9.5, color=DIMWHITE, align=PP_ALIGN.CENTER)
    if i < len(flow_items) - 1:
        add_textbox(s2, "→", fx + box_w, flow_y + Inches(0.52), arrow_w, Inches(0.4),
                    font_size=18, color=DIMWHITE, align=PP_ALIGN.CENTER)
    fx += box_w + arrow_w

# Two detail cards
card_y = Inches(3.05)
card_h = Inches(3.6)
card_w = Inches(5.8)

for cx, title, color, bullets in [
    (Inches(0.6), "Single Message Analysis", BLUE,
     ["Classifies into one of 4 labels",
      "Confidence scores for all classes",
      "Highlights matched toxic phrases",
      "Sanity check + response scripts"]),
    (Inches(6.9), "Conversation Window Audit", PURPLE,
     ["Analyzes last 5 messages in context",
      "Detects escalating toxic patterns",
      "Dominant pattern labeling across turns",
      "Overall severity score"]),
]:
    add_box(s2, cx, card_y, card_w, card_h,
            fill_color=RGBColor(0x0c, 0x14, 0x24),
            line_color=color)
    add_textbox(s2, title, cx + Inches(0.15), card_y + Inches(0.15),
                card_w - Inches(0.3), Inches(0.42),
                font_size=14, bold=True, color=color)
    body = "\n".join(f"• {b}" for b in bullets)
    add_textbox(s2, body, cx + Inches(0.15), card_y + Inches(0.65),
                card_w - Inches(0.3), card_h - Inches(0.9),
                font_size=12, color=DIMWHITE)

# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — Evaluation Plan
# ════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
set_bg(s3, SLIDE_BGS[2])

add_textbox(s3, "Evaluation Plan",
            Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
            font_size=36, bold=True, color=WHITE)
add_textbox(s3, "03 / 05", Inches(11.8), Inches(0.15), Inches(1.3), Inches(0.35),
            font_size=10, color=DIMWHITE, align=PP_ALIGN.RIGHT)

eval_cards = [
    (GREEN,  "① Sarcastic Sentiment Flip Test",
     '10 "positive-worded" insults — phrases that sound nice but mean the opposite.\n'
     'Standard classifiers label these Sincere. Ours should not.\n\n'
     'Pass threshold: ≥ 70% correctly identified as Sarcastic\n\n'
     'Example: "Oh wow, you\'re absolutely brilliant for doing that."'),
    (BLUE,   "② Contextual Consistency Audit",
     "A single sarcastic remark may be harmless.\n"
     "Repeated patterns across 5 messages reveal systemic abuse.\n\n"
     "• Counts toxic messages in conversation window\n"
     "• Flags escalation when ≥ 2 of last 3 are toxic\n"
     "• Dominant pattern detection across turns"),
    (AMBER,  "③ Held-out Test Set — F1 / Precision / Recall",
     "Per-class classification report on curated labeled examples.\n"
     "Confusion matrix reveals which classes the model conflates.\n\n"
     "Labels: Sincere · Sarcastic · Passive-Aggressive · Gaslighting"),
    (PURPLE, "④ Stress Tests (slang, emoji, fragments)",
     "Custom edge-case inputs: incomplete sentences, mixed-language,\n"
     "emoji substitutions, run-on text.\n\n"
     "Surfaces robustness failures before deployment."),
]

cw = Inches(5.9)
ch = Inches(2.85)
positions = [
    (Inches(0.35), Inches(1.3)),
    (Inches(6.7),  Inches(1.3)),
    (Inches(0.35), Inches(4.35)),
    (Inches(6.7),  Inches(4.35)),
]
for (ex, ey), (color, title, body) in zip(positions, eval_cards):
    add_box(s3, ex, ey, cw, ch,
            fill_color=RGBColor(0x08, 0x18, 0x0e),
            line_color=color)
    add_textbox(s3, title, ex + Inches(0.15), ey + Inches(0.12),
                cw - Inches(0.3), Inches(0.42),
                font_size=13, bold=True, color=color)
    add_textbox(s3, body, ex + Inches(0.15), ey + Inches(0.6),
                cw - Inches(0.3), ch - Inches(0.7),
                font_size=10.5, color=DIMWHITE)

# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — Limitations & Improvement
# ════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
set_bg(s4, SLIDE_BGS[3])

add_textbox(s4, "What Evaluation Revealed",
            Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
            font_size=36, bold=True, color=WHITE)
add_textbox(s4, "04 / 05", Inches(11.8), Inches(0.15), Inches(1.3), Inches(0.35),
            font_size=10, color=DIMWHITE, align=PP_ALIGN.RIGHT)

limit_cards = [
    (RED,    "Failure: Zero-shot misses subtle sarcasm",
     'bart-large-mnli scored "Oh wow, just wow" as Sincere (68% conf).\nFlip-test accuracy: 54% — below 70% threshold.'),
    (AMBER,  "Fix: Hybrid rule boost",
     "Added 30% regex pattern weight for zero-shot mode.\nPushed sarcasm flip accuracy to 78%."),
    (GREEN,  "Fix: Fine-tuning RoBERTa",
     "Fine-tuned on curated data; reduced rule weight to 15%.\nModel is now more reliable; further lifts F1 on PA class."),
    (RED,    "Failure: Single-message blind spots",
     '"Whatever, fine." alone = low-confidence PA.\n3rd time in 5 messages → escalation flag triggers.'),
    (AMBER,  "Insight: Confusion matrix",
     "Gaslighting ↔ Passive-Aggressive most confused pair.\nGuided us to add discriminating phrases to rule library."),
    (BLUE,   "Insight: Context window matters",
     "Individual accuracy ≠ user safety.\nConversation-level audit catches what single-turn misses."),
]

cw3 = Inches(3.9)
ch3 = Inches(2.45)
cols = [Inches(0.35), Inches(4.65), Inches(8.95)]
rows = [Inches(1.25), Inches(3.9)]

for idx, (color, title, body) in enumerate(limit_cards):
    ex = cols[idx % 3]
    ey = rows[idx // 3]
    add_box(s4, ex, ey, cw3, ch3,
            fill_color=RGBColor(0x1a, 0x10, 0x02),
            line_color=color)
    add_textbox(s4, title, ex + Inches(0.14), ey + Inches(0.12),
                cw3 - Inches(0.28), Inches(0.42),
                font_size=12, bold=True, color=color)
    add_textbox(s4, body, ex + Inches(0.14), ey + Inches(0.6),
                cw3 - Inches(0.28), ch3 - Inches(0.65),
                font_size=10.5, color=DIMWHITE)

# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — Demo & Impact
# ════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
set_bg(s5, SLIDE_BGS[4])

add_textbox(s5, "🛡️  Gaslight Guard — Demo & Impact",
            Inches(0.6), Inches(0.3), Inches(12.1), Inches(0.8),
            font_size=34, bold=True, color=PURPLE)
add_textbox(s5, "05 / 05", Inches(11.8), Inches(0.15), Inches(1.3), Inches(0.35),
            font_size=10, color=DIMWHITE, align=PP_ALIGN.RIGHT)

# Two cards
for cx, color, title, bullets in [
    (Inches(0.4), PURPLE, "What the app does",
     ["Paste any message → instant label + severity",
      "Conversation mode: analyze 5-turn chat history",
      'Sanity check: "You are NOT being overly sensitive"',
      "Coping strategies + suggested response scripts",
      "Direct links to crisis helplines (24/7)"]),
    (Inches(6.8), GREEN, "Who it helps",
     ["Students in unhealthy relationships",
      "Counselors triaging messages at scale",
      "Anyone who feels something is wrong but can't name it",
      "",
      "Resources: Crisis Text Line · National DV Hotline · SAMHSA · 7 Cups"]),
]:
    add_box(s5, cx, Inches(1.25), Inches(5.9), Inches(3.5),
            fill_color=RGBColor(0x12, 0x06, 0x22),
            line_color=color)
    add_textbox(s5, title, cx + Inches(0.15), Inches(1.38),
                Inches(5.6), Inches(0.42),
                font_size=14, bold=True, color=color)
    add_textbox(s5, "\n".join(f"• {b}" if b else "" for b in bullets),
                cx + Inches(0.15), Inches(1.88),
                Inches(5.6), Inches(2.6),
                font_size=11.5, color=DIMWHITE)

# Stats row
stats5 = [
    ("78%",    "sarcasm flip accuracy\nafter rule boost",      PURPLE),
    ("4",      "complementary\nevaluation axes",               GREEN),
    ("≤ 5s",   "real-time analysis\nper conversation",         AMBER),
]
sw = Inches(3.2)
sx5 = Inches(0.55)
for val, lbl, color in stats5:
    add_textbox(s5, val, sx5, Inches(5.1), sw, Inches(0.75),
                font_size=38, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(s5, lbl, sx5, Inches(5.9), sw, Inches(0.6),
                font_size=11, color=DIMWHITE, align=PP_ALIGN.CENTER)
    sx5 += Inches(3.4)

add_textbox(s5,
            "This tool is for educational purposes only and does not replace professional mental health support.",
            Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.4),
            font_size=9, color=DIMWHITE, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
out = "/home/user/540_hackathon_2/NLP product/frontend/public/Gaslight_Guard_Pitch.pptx"
prs.save(out)
print(f"Saved: {out}")
